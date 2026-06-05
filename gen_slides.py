from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0b, 0x0e, 0x1a)
INK = RGBColor(0xec, 0xed, 0xf6)
SOFT = RGBColor(0xa7, 0xad, 0xc8)
GOOD = RGBColor(0x67, 0xe8, 0xc5)
GOOD2 = RGBColor(0x9a, 0xf2, 0xd3)
TOXIC = RGBColor(0xc9, 0x3b, 0x8a)
ACCENT = RGBColor(0xf0, 0xa0, 0x5a)
ACCENT2 = RGBColor(0xfb, 0xd4, 0x9a)
DANGER = RGBColor(0xff, 0x5d, 0x6c)
GOLD = RGBColor(0xf6, 0xc6, 0x4a)
LINE = RGBColor(0x23, 0x28, 0x44)
CARD_BG = RGBColor(0x0e, 0x12, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_para(tf, text, size=18, color=SOFT, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri', spacing_after=Pt(4)):
    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '' else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = spacing_after
    return p

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

def add_tag(slide, left, top, text, color=ACCENT):
    tag_w = Inches(3.2)
    tag_h = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, tag_w, tag_h)
    shape.fill.background()
    r = min(int(str(color)[0:2], 16) + 40, 255)
    g = min(int(str(color)[2:4], 16) + 40, 255)
    b = min(int(str(color)[4:6], 16) + 40, 255)
    shape.line.color.rgb = RGBColor(r, g, b)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    set_para(tf, text, size=9, color=color, bold=True, align=PP_ALIGN.CENTER, font_name='Consolas')

def add_bullet_card(slide, left, top, width, height, title, items, title_color=GOOD2):
    card = add_card(slide, left, top, width, height)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    set_para(tf, title, size=14, color=title_color, bold=True)
    for item in items:
        set_para(tf, f"• {item}", size=11, color=SOFT)

# ==================== SLIDE 1: CAPA ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(4.8), Inches(1.0), "ENGENHARIA DE SOFTWARE EDUCATIVO · 2026.1", ACCENT)

tb = add_textbox(s, Inches(1), Inches(1.8), Inches(11.3), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
set_para(tf, "Eco Digital", size=48, color=INK, bold=True, align=PP_ALIGN.CENTER)
p2 = tf.add_paragraph()
p2.text = "O Labirinto da Reputação"
p2.font.size = Pt(36)
p2.font.bold = True
p2.font.color.rgb = GOOD
p2.font.name = 'Calibri'
p2.alignment = PP_ALIGN.CENTER

tb = add_textbox(s, Inches(2), Inches(4.2), Inches(9.3), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
set_para(tf, "Protótipo de jogo educativo sobre cyberbullying para adolescentes de 12 a 17 anos", size=16, color=SOFT, align=PP_ALIGN.CENTER)

tb = add_textbox(s, Inches(2), Inches(5.2), Inches(9.3), Inches(0.6))
tf = tb.text_frame
tf.word_wrap = True
set_para(tf, "Arthur Lopes  ·  Marcos Pierre  ·  João Lucas  ·  Jeová Bezerra", size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)

tb = add_textbox(s, Inches(2), Inches(6.2), Inches(9.3), Inches(0.4))
tf = tb.text_frame
set_para(tf, "05 de junho de 2026", size=12, color=SOFT, align=PP_ALIGN.CENTER)

# ==================== SLIDE 2: O PROBLEMA ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "O PROBLEMA", TOXIC)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Cyberbullying no Brasil", size=36, color=INK, bold=True)

# Stat cards
card = add_card(s, Inches(0.6), Inches(2.0), Inches(3.5), Inches(2.0))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
set_para(tf, "37%", size=52, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "dos adolescentes brasileiros já sofreram cyberbullying", size=12, color=SOFT, align=PP_ALIGN.CENTER)
set_para(tf, "Fonte: UNICEF", size=9, color=SOFT, align=PP_ALIGN.CENTER)

card = add_card(s, Inches(4.4), Inches(2.0), Inches(3.5), Inches(2.0))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
set_para(tf, "1 em 5", size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "crianças já passou por violência digital", size=12, color=SOFT, align=PP_ALIGN.CENTER)
set_para(tf, "Fonte: TIC Kids Online Brasil", size=9, color=SOFT, align=PP_ALIGN.CENTER)

add_bullet_card(s, Inches(8.3), Inches(2.0), Inches(4.4), Inches(2.0), "Impactos", [
    "Ansiedade e depressão",
    "Queda de autoestima",
    "Abandono escolar",
    "Risco de automutilação",
], title_color=DANGER)

card = add_card(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.4))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
set_para(tf, "Agravante digital", size=16, color=DANGER, bold=True)
set_para(tf, "Atinge a vítima 24h por dia, se espalha rápido e o agressor pode agir sob anonimato. No Brasil, a Lei 14.811/2024 tornou hediondo o crime de indução ao suicídio contra menores.", size=13, color=SOFT)

# ==================== SLIDE 3: O QUE É O ECO DIGITAL ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "VISÃO GERAL", GOOD)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "O que é o Eco Digital?", size=36, color=INK, bold=True)

tb = add_textbox(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.6))
tf = tb.text_frame
set_para(tf, "Jogo educativo de plataforma 2D + puzzle, single-player, baseado em navegador. Zero dependências — basta abrir o HTML.", size=15, color=SOFT)

add_bullet_card(s, Inches(0.6), Inches(2.6), Inches(3.8), Inches(2.4), "🌍  3 Mundos", [
    "Mundo 1: Feed (rede social)",
    "Mundo 2: Chat (mensageiro)",
    "Mundo 3: Fórum anônimo",
    "Cada mundo: 3 fases + 1 chefão",
])

add_bullet_card(s, Inches(4.8), Inches(2.6), Inches(3.8), Inches(2.4), "🎮  12 Fases", [
    "Progressão linear",
    "Dificuldade crescente",
    "Chefões com mecânica própria",
    "Medalhas (bronze/prata/ouro)",
])

add_bullet_card(s, Inches(9.0), Inches(2.6), Inches(3.8), Inches(2.4), "🧩  6 Tipos de Puzzle", [
    "A — Rede de Apoio",
    "B — Frase (ataque vs. fato)",
    "C — Identificar a Lei",
    "D — Ordem da Denúncia",
    "E — Empatia em Diálogo",
    "F — Reconstruir Mensagem",
])

# Flow
card = add_card(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.0))
tf = card.text_frame
tf.word_wrap = True
set_para(tf, "Reconhecer  →  Documentar  →  Agir  →  Denunciar  →  Restaurar", size=20, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "O jogador segue essa jornada em cada mundo, aprendendo que o enfrentamento ao cyberbullying é progressivo.", size=12, color=SOFT, align=PP_ALIGN.CENTER)

# ==================== SLIDE 4: CÍRCULO MÁGICO ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "FUNDAMENTAÇÃO TEÓRICA", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Círculo Mágico", size=36, color=INK, bold=True)

tb = add_textbox(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5))
tf = tb.text_frame
set_para(tf, "Conceito de Johan Huizinga (1938) — espaço temporário onde as regras do mundo real são suspensas e as regras do jogo prevalecem.", size=14, color=SOFT)

# Left card - mundo real
card = add_card(s, Inches(0.6), Inches(2.5), Inches(5.8), Inches(3.6))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
set_para(tf, "Mundo Real", size=18, color=DANGER, bold=True)
set_para(tf, "", size=6)
set_para(tf, "Cyberbullying → Medo, vergonha, isolamento", size=13, color=SOFT)
set_para(tf, "Consequência → Dano emocional permanente", size=13, color=SOFT)
set_para(tf, "Papel → Vítima passiva ou espectador omisso", size=13, color=SOFT)
set_para(tf, "Saída → Não há saída fácil", size=13, color=SOFT)

# Right card - mundo do jogo
card = add_card(s, Inches(6.8), Inches(2.5), Inches(5.8), Inches(3.6))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
set_para(tf, "Mundo do Jogo (Círculo Mágico)", size=18, color=GOOD, bold=True)
set_para(tf, "", size=6)
set_para(tf, "Conflito → Nuvens de Toxicidade (desafios superáveis)", size=13, color=SOFT)
set_para(tf, "Consequência → Dano à Resiliência (recuperável, retry)", size=13, color=SOFT)
set_para(tf, "Papel → Agente ativo que denuncia e resolve", size=13, color=SOFT)
set_para(tf, "Saída → Completar a fase = restaurar o ambiente", size=13, color=SOFT)

tb = add_textbox(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.6))
tf = tb.text_frame
set_para(tf, "O adolescente experimenta o cyberbullying de forma segura — sente a pressão, mas dentro de limites controláveis. A empatia com a situação é real; o risco é simulado.", size=13, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)

# ==================== SLIDE 5: MECÂNICAS ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "FUNDAMENTAÇÃO TEÓRICA", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Mecânicas do Jogo", size=36, color=INK, bold=True)

tb = add_textbox(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5))
tf = tb.text_frame
set_para(tf, "Regras e sistemas com os quais o jogador interage (Hunicke, LeBlanc & Zubek, 2004)", size=13, color=SOFT)

mechanics = [
    ("1", "Resiliência Emocional", "Barra de vida. Cyberbullying consome saúde mental; proteger-se é estratégia, não fraqueza."),
    ("2", "Nuvens → Plataformas Solidárias", "Resolver puzzle transforma toxidade em espaço seguro. Ignorar não resolve."),
    ("3", "6 Puzzles de Reparação", "Rede de apoio, frase, lei, ordem, empatia e reconstrução — cada um ensina algo específico."),
    ("4", "Prints de Segurança", "Colecionáveis que desbloqueiam conteúdo educativo real (leis, canais, estatísticas)."),
    ("5", "Progressão Linear", "Reconhecer → agir → denunciar. Não se pula etapas no enfrentamento."),
    ("6", "Medalhas e Rejogabilidade", "Bronze (completou), Prata (protegeu resiliência), Ouro (coletou tudo). Incentiva aprofundar."),
]

y = Inches(2.3)
for num, title, desc in mechanics:
    # Number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = GOOD
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.word_wrap = False
    set_para(ctf, num, size=14, color=RGBColor(0x0a, 0x1a, 0x15), bold=True, align=PP_ALIGN.CENTER)
    ctf.paragraphs[0].space_before = Pt(0)
    ctf.paragraphs[0].space_after = Pt(0)

    tb = add_textbox(s, Inches(1.5), y - Inches(0.05), Inches(11), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf, title, size=15, color=INK, bold=True)
    set_para(tf, desc, size=12, color=SOFT, spacing_after=Pt(2))
    y += Inches(0.78)

# ==================== SLIDE 6: JOGABILIDADE EDUCACIONAL ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "FUNDAMENTAÇÃO TEÓRICA", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Jogabilidade Educacional", size=36, color=INK, bold=True)

tb = add_textbox(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5))
tf = tb.text_frame
set_para(tf, "O conteúdo não é separado da jogabilidade — o jogador aprende jogando, não lendo (Prensky, 2001; Gee, 2003)", size=14, color=ACCENT2, bold=True)

# Table
rows = [
    ("Mecânica", "O que o jogador FAZ", "O que APRENDE"),
    ("Resiliência", "Gerencia recurso, evita dano", "Cyberbullying consome saúde mental"),
    ("Nuvens → Solidárias", "Resolve puzzle para transformar", "Enfrentar com conhecimento resolve"),
    ("Puzzle A — Rede de Apoio", "Ordena a cadeia de apoio", "Existe sequência de proteção"),
    ("Puzzle B — Frase", "Classifica ataque vs. fato", "Crítica ≠ ataque"),
    ("Puzzle C — Lei", "Escolhe a lei aplicável", "Há leis que protegem"),
    ("Puzzle D — Denúncia", "Ordena os passos", "Há protocolo a seguir"),
    ("Puzzle E — Empatia", "Escolhe resposta empática", "Escutar e validar > conselhos"),
    ("Puzzle F — Reconstruir", "Monta mensagem de apoio", "Palavras certas constroem"),
    ("Prints", "Coleta provas pelo nível", "Print é prova documental"),
]

table_shape = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.2))
table = table_shape.table
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(4.5)
table.columns[2].width = Inches(4.6)

for i, (c1, c2, c3) in enumerate(rows):
    for j, val in enumerate([c1, c2, c3]):
        cell = table.cell(i, j)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11) if i > 0 else Pt(11)
        p.font.bold = (i == 0)
        p.font.color.rgb = GOOD2 if i == 0 else SOFT
        p.font.name = 'Calibri'
        cell.fill.solid()
        if i == 0:
            cell.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x36)
        elif i % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(0x0e, 0x12, 0x28)
        else:
            cell.fill.fore_color.rgb = RGBColor(0x0b, 0x0f, 0x22)

# ==================== SLIDE 7: TELAS DO PROTÓTIPO ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "PROTÓTIPO", GOOD)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Telas do Jogo", size=36, color=INK, bold=True)

screens = [
    ("🏠", "Tela Inicial", "Título, descrição e botões\nJogar, Como Jogar, Sobre"),
    ("🗺️", "Mapa-Múndi", "3 mundos com 4 fases cada.\nNavegação por nós."),
    ("⚔️", "Gameplay", "Canvas 2D: plataforma, pulo,\ninimigos, nuvens, power-ups"),
    ("🧩", "Painel de Reparação", "6 tipos de puzzle educativo\ninterativos com verificação"),
    ("🏆", "Resultado", "Medalhas (bronze/prata/ouro)\ne checklist de objetivos"),
    ("📁", "Acervo de Provas", "20 cards educativos:\nleis, dados, canais, histórias"),
]

positions = [
    (Inches(0.6), Inches(2.0)),
    (Inches(4.6), Inches(2.0)),
    (Inches(8.6), Inches(2.0)),
    (Inches(0.6), Inches(4.3)),
    (Inches(4.6), Inches(4.3)),
    (Inches(8.6), Inches(4.3)),
]

for (icon, title, desc), (x, y) in zip(screens, positions):
    card = add_card(s, x, y, Inches(3.6), Inches(1.8))
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    set_para(tf, f"{icon}  {title}", size=16, color=GOOD2, bold=True)
    set_para(tf, desc, size=12, color=SOFT)

# ==================== SLIDE 8: HUD ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "PROTÓTIPO", GOOD)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Interface durante o jogo (HUD)", size=36, color=INK, bold=True)

add_bullet_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(1.8),
    "Barra de Resiliência", [
        "Saúde mental do adolescente no jogo",
        "Dano: inimigos, toxidade, erros",
        "Zerar = esgotamento (restart)",
        "Proteger-se é estratégia, não fraqueza",
    ], title_color=GOOD2)

add_bullet_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.8),
    "Mini-mapa e Timer", [
        "Posição do jogador e nuvens no nível",
        "Timer em fases de viralização",
        "Simula urgência — fofocas se espalham",
    ], title_color=ACCENT2)

add_bullet_card(s, Inches(0.6), Inches(4.1), Inches(5.8), Inches(1.8),
    "Prints de Segurança", [
        "Contador de provas coletadas",
        "Cada print desbloqueia conteúdo no Acervo",
        "Print = prova documental (primeiro passo)",
    ], title_color=ACCENT2)

add_bullet_card(s, Inches(6.8), Inches(4.1), Inches(5.8), Inches(1.8),
    "Power-ups", [
        "Escudo de Empatia — absorve 1 ataque",
        "Bloqueio Temporário — congela inimigos 5s",
        "Aliado — atrai haters por 10s",
        "Print Dourado — vale 5 prints",
    ], title_color=GOOD2)

card = add_card(s, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7))
tf = card.text_frame
tf.word_wrap = True
set_para(tf, "Controles: ← → mover · Espaço pular · E interagir · M mudo · Esc voltar", size=13, color=SOFT, align=PP_ALIGN.CENTER, font_name='Consolas')

# ==================== SLIDE 9: CONTEÚDO EDUCATIVO ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "CONTEÚDO INTEGRADO", GOOD)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "O que o jogo ensina", size=36, color=INK, bold=True)

add_bullet_card(s, Inches(0.6), Inches(2.0), Inches(3.8), Inches(2.6),
    "Legislação Brasileira", [
        "ECA — proteção integral online",
        "Lei 13.185/2015 — bullying",
        "Lei 14.811/2024 — crime hediondo",
        "Lei Carolina Dieckmann — vazamento",
    ], title_color=GOOD2)

add_bullet_card(s, Inches(4.8), Inches(2.0), Inches(3.8), Inches(2.6),
    "Canais de Denúncia", [
        "SaferNet — denúncia anônima",
        "Disque 100 — Direitos Humanos 24h",
        "CVV 188 — apoio emocional 24h",
        "Polícia Federal — crimes cibernéticos",
    ], title_color=GOOD2)

add_bullet_card(s, Inches(9.0), Inches(2.0), Inches(3.8), Inches(2.6),
    "Habilidades Socioemocionais", [
        "Separar crítica de ataque",
        "Praticar empatia ativa",
        "Pedir ajuda como estratégia",
        "Não responder com ódio",
        "Acionar adulto não é 'X9'",
    ], title_color=GOOD2)

# Protocolo
card = add_card(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.8))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
set_para(tf, "Protocolo de Denúncia (Puzzle D)", size=16, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "", size=4)
set_para(tf, "① Tirar print  →  ② Bloquear o agressor  →  ③ Denunciar na plataforma  →  ④ Contar para um responsável  →  ⑤ Acionar canal externo (SaferNet, Disque 100, PF)", size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "", size=4)
set_para(tf, "O jogador pratica essa ordem múltiplas vezes ao longo do jogo — internaliza o protocolo pela repetição.", size=12, color=SOFT, align=PP_ALIGN.CENTER)

# ==================== SLIDE 10: DEMONSTRAÇÃO ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "AO VIVO", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Demonstração", size=36, color=INK, bold=True)

tb = add_textbox(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.5))
tf = tb.text_frame
set_para(tf, "Abrir o jogo no navegador e demonstrar o fluxo completo", size=15, color=ACCENT2, bold=True)

add_bullet_card(s, Inches(0.6), Inches(2.5), Inches(5.8), Inches(3.5),
    "Roteiro sugerido (~1 min)", [
        "1. Tela inicial → Como Jogar",
        "2. Mapa-Múndi → Selecionar fase 1-1",
        "3. Gameplay: andar, pular, coletar print",
        "4. Aproximar da Nuvem de Toxicidade",
        "5. Pressionar E → Abrir Painel de Reparação",
        "6. Resolver o puzzle",
        "7. Ver a nuvem virar Plataforma Solidária",
        "8. Acervo de Provas desbloqueado",
    ], title_color=GOOD2)

card = add_card(s, Inches(6.8), Inches(2.5), Inches(5.8), Inches(3.5))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
set_para(tf, "Pontos a destacar na demonstração", size=16, color=ACCENT2, bold=True)
set_para(tf, "", size=8)
set_para(tf, "→ A resiliência diminui perto de nuvens tóxicas (efeito sonoro + visual)", size=13, color=SOFT)
set_para(tf, "", size=4)
set_para(tf, "→ O puzzle ensina conteúdo real (classificar ataque vs. fato)", size=13, color=SOFT)
set_para(tf, "", size=4)
set_para(tf, "→ A nuvem resolvida vira caminho seguro (metáfora: denunciar transforma o espaço)", size=13, color=SOFT)
set_para(tf, "", size=4)
set_para(tf, "→ O print desbloqueia card no Acervo com lei, dado ou canal real", size=13, color=SOFT)
set_para(tf, "", size=4)
set_para(tf, "→ Chefão final: reunir provas + resolver puzzle = denúncia formal aceita", size=13, color=SOFT)

# ==================== SLIDE 11: ARTICULAÇÃO TEÓRICA ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "SÍNTESE", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Articulação dos Conceitos", size=36, color=INK, bold=True)

# Three boxes
card = add_card(s, Inches(0.6), Inches(2.0), Inches(3.6), Inches(2.8))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
set_para(tf, "Círculo Mágico", size=18, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "", size=4)
set_para(tf, "Espaço seguro onde o adolescente experimenta o cyberbullying sem consequências reais. A empatia é real; o risco é simulado.", size=12, color=SOFT, align=PP_ALIGN.CENTER)

# Arrow 1
tb = add_textbox(s, Inches(4.35), Inches(2.8), Inches(0.8), Inches(0.6))
tf = tb.text_frame
set_para(tf, "→", size=32, color=SOFT, bold=True, align=PP_ALIGN.CENTER)

card = add_card(s, Inches(5.2), Inches(2.0), Inches(3.0), Inches(2.8))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
set_para(tf, "Mecânicas", size=18, color=GOOD, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "", size=4)
set_para(tf, "Resiliência, puzzles, prints, progressão — regras que simulam a experiência real de enfrentamento ao cyberbullying.", size=12, color=SOFT, align=PP_ALIGN.CENTER)

# Arrow 2
tb = add_textbox(s, Inches(8.35), Inches(2.8), Inches(0.8), Inches(0.6))
tf = tb.text_frame
set_para(tf, "→", size=32, color=SOFT, bold=True, align=PP_ALIGN.CENTER)

card = add_card(s, Inches(9.2), Inches(2.0), Inches(3.6), Inches(2.8))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
set_para(tf, "Aprendizagem", size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "", size=4)
set_para(tf, "O jogador internaliza leis, canais, protocolos e empatia pela ação — não pela leitura.", size=12, color=SOFT, align=PP_ALIGN.CENTER)

# Bottom card
card = add_card(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.6))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
set_para(tf, "Quando o círculo mágico se fecha (tela de Vitória), o jogador leva para o mundo real:", size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)
set_para(tf, "", size=4)
set_para(tf, "Leis que protegem  ·  Canais que funcionam  ·  Protocolo que pode seguir  ·  Empatia como prática", size=16, color=GOOD2, bold=True, align=PP_ALIGN.CENTER)

# ==================== SLIDE 12: PRÓXIMOS PASSOS ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "ENCAMINHAMENTO", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Próximos Passos", size=36, color=INK, bold=True)

add_bullet_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(3.0),
    "Curto Prazo", [
        "Persistência com localStorage",
        "Suporte a touch/mobile",
        "Testes com público-alvo (12-17 anos)",
        "Acessibilidade (ARIA labels, contraste)",
        "Mais variação de puzzles",
    ], title_color=GOOD2)

add_bullet_card(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(3.0),
    "Médio Prazo", [
        "Novos mundos e tipos de puzzle",
        "Dashboard do professor",
        "Localização (inglês, espanhol)",
        "Avaliação de impacto educacional",
        "Deploy como PWA (funciona offline)",
    ], title_color=ACCENT2)

card = add_card(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.2))
tf = card.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
set_para(tf, "Tecnologia", size=16, color=GOOD2, bold=True)
set_para(tf, "Arquivo HTML único (~3000 linhas), zero dependências externas, Canvas 2D nativo, Web Audio API sintetizado. Funciona em qualquer navegador, sem instalação, sem internet após download.", size=13, color=SOFT)

# ==================== SLIDE 13: REFERÊNCIAS ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_tag(s, Inches(0.6), Inches(0.4), "REFERÊNCIAS", ACCENT)

tb = add_textbox(s, Inches(0.6), Inches(1.0), Inches(12), Inches(0.8))
tf = tb.text_frame
set_para(tf, "Referências", size=36, color=INK, bold=True)

refs_left = [
    "HUIZINGA, J. Homo Ludens. 1938.",
    "SALEN, K.; ZIMMERMAN, E. Rules of Play. MIT Press, 2004.",
    "HUNICKE, R. et al. MDA: A Formal Approach to Game Design. AAAI, 2004.",
    "SICART, M. Defining Game Mechanics. Game Studies, 2008.",
]
refs_right = [
    "GEE, J. P. What Video Games Have to Teach Us. 2003.",
    "PRENSKY, M. Digital Game-Based Learning. 2001.",
    "UNICEF. Estado Mundial da Infância 2017.",
    "Brasil. Lei 13.185/2015; Lei 14.811/2024; ECA (Lei 8.069/1990).",
]

add_bullet_card(s, Inches(0.6), Inches(1.8), Inches(5.8), Inches(3.0),
    "Teoria", refs_left, title_color=GOOD2)

add_bullet_card(s, Inches(6.8), Inches(1.8), Inches(5.8), Inches(3.0),
    "Prática e Legislação", refs_right, title_color=GOOD2)

# ==================== SLIDE 14: ENCERRAMENTO ====================
s = prs.slides.add_slide(BLANK)
set_bg(s)

tb = add_textbox(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2))
tf = tb.text_frame
set_para(tf, "Obrigado!", size=52, color=INK, bold=True, align=PP_ALIGN.CENTER)

tb = add_textbox(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6))
tf = tb.text_frame
set_para(tf, "Eco Digital: O Labirinto da Reputação", size=22, color=GOOD, bold=True, align=PP_ALIGN.CENTER)

tb = add_textbox(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6))
tf = tb.text_frame
set_para(tf, "Arthur Lopes  ·  Marcos Pierre  ·  João Lucas  ·  Jeová Bezerra", size=16, color=INK, bold=True, align=PP_ALIGN.CENTER)

tb = add_textbox(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4))
tf = tb.text_frame
set_para(tf, "Engenharia de Software Educativo · 2026.1", size=13, color=SOFT, align=PP_ALIGN.CENTER)

tb = add_textbox(s, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4))
tf = tb.text_frame
set_para(tf, "05 de junho de 2026", size=12, color=SOFT, align=PP_ALIGN.CENTER)

prs.save('Eco Digital - Apresentacao.pptx')
print("PPTX gerado com sucesso!")
